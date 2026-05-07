import os
os.environ["HF_HUB_DISABLE_PROGRESS_BARS"] = "1"
os.environ["TRANSFORMERS_NO_ADVISORY_WARNINGS"] = "1"
import fitz
import docx
from pptx import Presentation
from sentence_transformers import SentenceTransformer, CrossEncoder
import faiss
import numpy as np
import re
from typing import List, Tuple

# ==============================
# MODEL MANAGEMENT (CACHED)
# ==============================
_model = None
_cross_model = None

def get_model():
    global _model
    if _model is None:
        _model = SentenceTransformer("BAAI/bge-base-en-v1.5")
    return _model

def get_cross_encoder():
    global _cross_model
    if _cross_model is None:
        _cross_model = CrossEncoder("cross-encoder/ms-marco-MiniLM-L-6-v2")
    return _cross_model


# ==============================
# TEXT CLEANING & CHUNKING
# ==============================
def clean_text(text: str) -> str:
    return re.sub(r"\s+", " ", text).strip()

def chunk_text(text: str, chunk_size=5, overlap=2) -> List[str]:
    sentences = re.split(r'(?<=[.!?]) +', text)
    chunks = []

    for i in range(0, len(sentences), chunk_size - overlap):
        chunk = " ".join(sentences[i:i + chunk_size])
        if len(chunk.strip()) > 50:
            chunks.append(chunk)

    return chunks


# ==============================
# DOCUMENT EXTRACTION
# ==============================
def extract_pdf(path):
    try:
        doc = fitz.open(path)
        return clean_text(" ".join([p.get_text() for p in doc]))
    except:
        return ""

def extract_docx(path):
    try:
        doc = docx.Document(path)
        return clean_text("\n".join([p.text for p in doc.paragraphs]))
    except:
        return ""

def extract_pptx(path):
    try:
        prs = Presentation(path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return clean_text("\n".join(text))
    except:
        return ""


# ==============================
# LOAD DOCUMENTS
# ==============================
def load_documents(folder: str) -> Tuple[List[str], List[str]]:
    docs, names = [], []

    for file in os.listdir(folder):
        path = os.path.join(folder, file)

        if file.endswith(".pdf"):
            text = extract_pdf(path)
        elif file.endswith(".docx"):
            text = extract_docx(path)
        elif file.endswith(".pptx"):
            text = extract_pptx(path)
        else:
            continue

        if not text:
            continue

        for chunk in chunk_text(text):
            docs.append(chunk)
            names.append(file)

    return docs, names


# ==============================
# EMBEDDINGS (BATCHED)
# ==============================
def create_embeddings(texts: List[str]) -> np.ndarray:
    model = get_model()

    embeddings = model.encode(
        texts,
        batch_size=64,               # 🔥 batched for speed
        convert_to_numpy=True,
        show_progress_bar=False
    )

    faiss.normalize_L2(embeddings)
    return embeddings


# ==============================
# FAISS INDEX
# ==============================
def build_index(embeddings: np.ndarray):
    dim = embeddings.shape[1]

    # 🔥 scalable index (switch to IVF later easily)
    index = faiss.IndexFlatIP(dim)
    index.add(embeddings)

    return index


# ==============================
# UNIFIED SEARCH (DOC + IMAGE CAPTIONS)
# ==============================
def search(query, docs, names, index, top_k=5):
    model = get_model()
    cross = get_cross_encoder()

    query = "Represent this sentence for searching relevant documents: " + query

    # -------- RETRIEVAL -------- #
    q_vec = model.encode([query], convert_to_numpy=True)
    faiss.normalize_L2(q_vec)

    scores, indices = index.search(q_vec, top_k * 10)

    candidates = [
        (names[i], docs[i],
         "image" if names[i].lower().endswith((".jpg", ".png", ".jpeg")) else "doc")
        for i in indices[0] if i != -1
    ]

    # -------- RERANK -------- #
    pairs = [(query, doc[:300]) for _, doc, _ in candidates]
    rerank_scores = cross.predict(pairs)

    ranked = sorted(
        zip(candidates, rerank_scores),
        key=lambda x: x[1],
        reverse=True
    )

    # -------- DEDUP -------- #
    results, seen = [], set()

    for (name, doc, source), score in ranked:
        if name not in seen:
            results.append((name, doc, source, float(score)))
            seen.add(name)

        if len(results) >= top_k:
            break

    return results


# ==============================
# BEST SENTENCE EXTRACTION
# ==============================
def get_best_sentence(text, query):
    model = get_model()

    sentences = re.split(r'(?<=[.!?]) +', text)
    if len(sentences) <= 1:
        return text

    sent_emb = model.encode(sentences, convert_to_numpy=True)
    query_emb = model.encode([query], convert_to_numpy=True)

    faiss.normalize_L2(sent_emb)
    faiss.normalize_L2(query_emb)

    scores = np.dot(sent_emb, query_emb.T).squeeze()
    return sentences[np.argmax(scores)]